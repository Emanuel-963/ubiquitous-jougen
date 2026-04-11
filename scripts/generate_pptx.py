"""Generate a simple PPTX from `docs/PRESENTATION.md`.

Usage:
    python scripts/generate_pptx.py --in docs/PRESENTATION.md \
        --out docs/presentation.pptx

This script creates one slide per top-level section.
Smaller headings become bullets.
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Pt


def parse_markdown(md_text: str) -> List[dict]:
    """Parse very small subset of markdown into sections.

    Returns list of {title: str, bullets: List[str]}
    """
    lines = md_text.splitlines()
    sections = []
    current = None
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            if current:
                sections.append(current)
            current = {"title": line[2:].strip(), "bullets": []}
        elif line.startswith("## "):
            if current is None:
                current = {"title": "", "bullets": []}
            current["bullets"].append(line[3:].strip())
        elif line.startswith("- ") or line.startswith("* "):
            if current is None:
                current = {"title": "", "bullets": []}
            current["bullets"].append(line[2:].strip())
        else:
            # treat as paragraph -> bullet
            if current is None:
                current = {"title": "", "bullets": []}
            current["bullets"].append(line)
    if current:
        sections.append(current)
    return sections


def make_pptx(sections: List[dict], out_path: Path) -> None:
    prs = Presentation()
    # Title slide
    if sections:
        first = sections[0]
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = first.get("title", "Presentation")
        subtitle = slide.placeholders[1]
        subtitle.text = "Generated from docs/PRESENTATION.md"

    for sec in sections[1:]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = sec.get("title", "")
        body = slide.shapes.placeholders[1].text_frame
        for b in sec.get("bullets", []):
            p = body.add_paragraph() if body.text else body.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main(argv=None):
    p = argparse.ArgumentParser()
    p.add_argument("--in", dest="md", default="docs/PRESENTATION.md")
    p.add_argument("--out", dest="out", default="docs/presentation.pptx")
    args = p.parse_args(argv)

    md_path = Path(args.md)
    if not md_path.exists():
        print("Input markdown not found:", md_path)
        return 2
    md = md_path.read_text(encoding="utf-8")
    sections = parse_markdown(md)
    make_pptx(sections, Path(args.out))
    print("Generated:", args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
